#!/usr/bin/env python3
"""Export presentation.md to PPTX for Google Slides import."""

import argparse
import os
import re
import subprocess
import sys
from pathlib import Path


def ensure_package(import_name: str, package: str, allow_install: bool) -> bool:
    try:
        __import__(import_name)
        return True
    except ImportError:
        if not allow_install:
            return False
        print(f"Missing dependency: {package}. Attempting install...")
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", package],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            print(result.stdout)
            print(result.stderr)
            return False
        try:
            __import__(import_name)
            return True
        except ImportError:
            return False


def parse_slides(markdown: str):
    sections = [s.strip() for s in markdown.split("---") if s.strip()]
    slides = []
    for section in sections:
        title = None
        bullets = []
        images = []
        for line in section.splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith("#"):
                if title is None:
                    title = line.lstrip("#").strip()
                else:
                    bullets.append(line.lstrip("#").strip())
                continue
            match = re.match(r"[-*]\s+(.*)", line)
            if match:
                bullets.append(match.group(1))
                continue
            img_match = re.match(r"!\[[^\]]*\]\(([^)]+)\)", line)
            if img_match:
                images.append(img_match.group(1))
                continue
        slides.append({"title": title, "bullets": bullets, "images": images})
    return slides


def main():
    parser = argparse.ArgumentParser(description="Export presentation.md to PPTX.")
    parser.add_argument(
        "--input",
        default="presentation.md",
        help="Markdown slide deck input.",
    )
    parser.add_argument(
        "--out",
        default="artifacts/presentation.pptx",
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--no-install",
        action="store_true",
        help="Disable auto-install of python-pptx.",
    )
    args = parser.parse_args()

    if not ensure_package("pptx", "python-pptx", not args.no_install):
        raise SystemExit("python-pptx is required")

    from pptx import Presentation
    from pptx.util import Inches

    base_dir = Path(__file__).resolve().parent
    input_path = base_dir / args.input
    out_path = base_dir / args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    markdown = input_path.read_text(encoding="utf-8")
    slides = parse_slides(markdown)

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for idx, slide_data in enumerate(slides):
        layout = title_layout if idx == 0 else content_layout
        slide = prs.slides.add_slide(layout)

        title = slide_data.get("title") or ""
        if slide.shapes.title:
            slide.shapes.title.text = title

        if layout == content_layout and slide.placeholders:
            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in slide_data.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0

        for image_path in slide_data.get("images", []):
            image_file = (base_dir / image_path).resolve()
            if not image_file.exists():
                image_file = (base_dir / "artifacts" / image_path).resolve()
            if image_file.exists():
                slide.shapes.add_picture(
                    str(image_file),
                    Inches(1),
                    Inches(1.5),
                    width=Inches(8),
                )

    prs.save(out_path)
    print(f"Wrote slides to {out_path}")


if __name__ == "__main__":
    main()
